from flask import Flask, request, send_file, jsonify
from flask_cors import CORS
from dotenv import load_dotenv
load_dotenv()

from io import BytesIO
from PyPDF2 import PdfReader
from fpdf import FPDF
import os
import time
from functools import wraps
import requests
from docx import Document
from pptx import Presentation

app = Flask(__name__)
CORS(app)

# ---------------- RATE LIMIT ----------------
request_count = {}
RATE_LIMIT = 50

def rate_limit(f):
    @wraps(f)
    def decorated_function(*args, **kwargs):
        ip = request.remote_addr
        current_time = time.time()

        if ip in request_count:
            request_count[ip] = [t for t in request_count[ip] if current_time - t < 3600]
        else:
            request_count[ip] = []

        if len(request_count[ip]) >= RATE_LIMIT:
            return "Rate limit exceeded", 429

        request_count[ip].append(current_time)
        return f(*args, **kwargs)
    return decorated_function

# ---------------- API KEY ----------------
GROQ_API_KEY = os.getenv("GROQ_API_KEY")

# ---------------- FILE EXTRACTION ----------------
def extract_text_from_pdf(file_path):
    reader = PdfReader(file_path)
    text = ""
    for page in reader.pages:
        t = page.extract_text()
        if t:
            text += t + "\n"
    return text

def extract_text_from_docx(file_path):
    doc = Document(file_path)
    return "\n".join([p.text for p in doc.paragraphs])

def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def extract_text_from_file(path, filename):
    ext = filename.split('.')[-1].lower()
    if ext == "pdf":
        return extract_text_from_pdf(path)
    elif ext == "docx":
        return extract_text_from_docx(path)
    elif ext == "pptx":
        return extract_text_from_pptx(path)
    elif ext == "txt":
        return open(path).read()
    return ""

# ---------------- GROQ ----------------
def generate_with_groq(prompt, query):
    try:
        if not GROQ_API_KEY:
            return "❌ GROQ API KEY NOT SET"

        response = requests.post(
            "https://api.groq.com/openai/v1/chat/completions",
            headers={
                "Authorization": f"Bearer {GROQ_API_KEY}",
                "Content-Type": "application/json"
            },
            json={
                "model": "llama3-70b-8192",  # ✅ FIXED MODEL
                "messages": [
                    {"role": "system", "content": prompt},
                    {"role": "user", "content": query}
                ],
                "temperature": 0.7
            }
        )

        if response.status_code == 200:
            data = response.json()
            print("API RESPONSE:", data)
            return data.get("choices", [{}])[0].get("message", {}).get("content", "No output")

        else:
            print("ERROR:", response.text)
            return "Error generating content"

    except Exception as e:
        print("EXCEPTION:", e)
        return "Error"

# ---------------- GENERATION ----------------
def generate_questions_with_answers(prompt, qtype, num):
    system = "You are an expert teacher. Generate questions with answers."
    query = f"Create {num} {qtype} questions with answers from:\n{prompt}"
    return generate_with_groq(system, query)

# ---------------- PDF ----------------
def save_questions_to_pdf(text):
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", size=12)

    # ✅ HANDLE ERROR CASE
    if "❌" in text or text.strip() == "":
        text = "Error: Could not generate questions. Please try again."

    for line in text.split("\n"):
        pdf.multi_cell(0, 8, line)

    # ✅ CORRECT BYTE GENERATION
    pdf_bytes = pdf.output(dest='S').encode('latin-1')

    output = BytesIO(pdf_bytes)
    output.seek(0)

    return output
# ---------------- ROUTES ----------------
@app.route('/')
def index():
    return "Backend Running ✅"

@app.route('/generate', methods=['POST'])
@rate_limit
def generate():
    try:
        file = request.files.get("pdf_file")

        if not file:
            return "No file uploaded", 400

        filename = file.filename
        temp_path = "temp_" + filename
        file.save(temp_path)

        text = extract_text_from_file(temp_path, filename)

        questions = generate_questions_with_answers(
            text,
            request.form.get("question_type", "MCQ"),
            int(request.form.get("num_questions", 5))
        )

        pdf = save_questions_to_pdf(questions)

        os.remove(temp_path)

        return send_file(pdf, download_name="questions.pdf", as_attachment=True)

    except Exception as e:
        print("ERROR IN /generate:", e)
        return "Server Error", 500

@app.route('/health')
def health():
    return jsonify({"status": "ok"})

# ---------------- RUN ----------------
if __name__ == "__main__":
    port = int(os.getenv("PORT", 5000))
    app.run(host="0.0.0.0", port=port)
